from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor as PptRGBColor
import os

# Colores Duoc
DUOC_BLUE = RGBColor(0, 102, 204)
DUOC_DARK_BLUE = RGBColor(0, 68, 153)

print("Creando documentos...")

# ==================== CREAR DOCUMENTO WORD ====================
print("\n1. Creando documento Word...")
doc = Document()

# Configurar márgenes
sections = doc.sections
for section in sections:
    section.top_margin = Inches(1)
    section.bottom_margin = Inches(1)
    section.left_margin = Inches(1)
    section.right_margin = Inches(1)

# PORTADA
heading = doc.add_heading('🎬 GUION PARA VIDEO', 0)
heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
heading.runs[0].font.color.rgb = DUOC_BLUE

subtitle = doc.add_heading('EVALUACIÓN FINAL TRANSVERSAL SEMANA 9', 2)
subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
subtitle.runs[0].font.color.rgb = DUOC_DARK_BLUE

# Información del encabezado
info = doc.add_paragraph()
info.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = info.add_run('Sistema Teatro Moro - Gestión de Ventas de Entradas\n')
run.font.size = Pt(14)
run.bold = True
run = info.add_run('Fundamentos de Programación (PRY2201)\n')
run.font.size = Pt(12)
run = info.add_run('Duración: 5 minutos máximo\n')
run.font.size = Pt(12)
run = info.add_run('Estudiante: Fuad Oñate - Programación Online')
run.font.size = Pt(12)
run.bold = True
run.font.color.rgb = DUOC_BLUE

doc.add_page_break()

# INTRODUCCIÓN
doc.add_heading('📍 INTRODUCCIÓN (30 segundos)', 1).runs[0].font.color.rgb = DUOC_BLUE
p = doc.add_paragraph()
p.add_run('Hola, mi nombre es Fuad Oñate y les presento el Sistema de Gestión de Ventas de Entradas para el Teatro Moro, desarrollado como Evaluación Final Transversal de la asignatura Fundamentos de Programación, semana 9.\n\n').font.size = Pt(11)
p.add_run('Este sistema permite vender entradas, gestionar asientos y aplicar descuentos automáticos según el tipo de cliente.').font.size = Pt(11)

# ARQUITECTURA
doc.add_heading('🏗️ ARQUITECTURA DEL SISTEMA (45 segundos)', 1).runs[0].font.color.rgb = DUOC_BLUE
p = doc.add_paragraph()
run = p.add_run('📺 MOSTRAR CÓDIGO EN NETBEANS - LÍNEAS 14-39\n\n')
run.font.size = Pt(10)
run.bold = True
run.font.color.rgb = RGBColor(255, 165, 0)

p.add_run('El sistema está construido con:\n\n').font.size = Pt(11)

bullet_items = [
    '10 arrays paralelos que almacenan información de 150 asientos, incluyendo datos del cliente, precios y estados.',
    'Un ArrayList para el historial dinámico de ventas.',
    'Variables estáticas que rastrean las ventas totales, ingresos y descuentos aplicados.',
    '5 constantes para los precios de cada ubicación: VIP a 60 mil pesos, Palco a 50 mil, Platea Baja a 40 mil, Platea Alta a 30 mil y Galería a 20 mil pesos.'
]

for item in bullet_items:
    p = doc.add_paragraph(item, style='List Bullet')
    p.runs[0].font.size = Pt(11)

p = doc.add_paragraph()
p.add_run('\nEl teatro tiene 150 asientos divididos en estas 5 secciones de 30 asientos cada una.').font.size = Pt(11)

# IMPLEMENTACIÓN
doc.add_heading('⚙️ IMPLEMENTACIÓN Y ESTRATEGIAS (1 minuto)', 1).runs[0].font.color.rgb = DUOC_BLUE
p = doc.add_paragraph()
run = p.add_run('📺 MOSTRAR MÉTODO venderEntrada()\n\n')
run.font.size = Pt(10)
run.bold = True
run.font.color.rgb = RGBColor(255, 165, 0)

p.add_run('Implementé el sistema usando múltiples estructuras de control:\n\n').font.size = Pt(11)

impl_items = [
    'Estructuras IF-ELSE para validar edad y aplicar descuentos automáticos. Si el cliente es menor de 12 años, aplica 5% de descuento. Si tiene 65 años o más, aplica 30% para tercera edad.',
    'Un SWITCH para el menú principal y la selección de ubicaciones.',
    'Ciclos FOR para recorrer los asientos y buscar disponibilidad.',
    'Un DO-WHILE para mantener el menú activo hasta que el usuario decida salir.'
]

for item in impl_items:
    p = doc.add_paragraph(item, style='List Bullet')
    p.runs[0].font.size = Pt(11)

p = doc.add_paragraph()
p.add_run('\nEl manejo de excepciones está implementado con bloques TRY-CATCH en todas las entradas del usuario, protegiendo el sistema de datos inválidos.').font.size = Pt(11)

# DIFICULTADES
doc.add_heading('🔧 DIFICULTADES Y SOLUCIONES (45 segundos)', 1).runs[0].font.color.rgb = DUOC_BLUE
p = doc.add_paragraph()
run = p.add_run('📺 MOSTRAR MÉTODO modificarAsiento() - LÍNEAS 304-317\n\n')
run.font.size = Pt(10)
run.bold = True
run.font.color.rgb = RGBColor(255, 165, 0)

p.add_run('Durante el desarrollo enfrenté varios desafíos:\n\n').font.size = Pt(11)
p.add_run('El principal fue la modificación de asientos cuando el cliente cambia de ubicación. La dificultad era recalcular el precio manteniendo el descuento proporcional.\n\n').font.size = Pt(11)
p.add_run('Lo resolví calculando primero el porcentaje de descuento original, luego aplicando ese mismo porcentaje al precio de la nueva ubicación, y finalmente ajustando los totales del sistema.\n\n').font.size = Pt(11)
p.add_run('Otro desafío fue la validación de RUT para evitar compras duplicadas. Lo solucioné recorriendo todos los asientos vendidos y comparando el RUT ingresado.').font.size = Pt(11)

doc.add_page_break()

# DEMOSTRACIÓN
doc.add_heading('💻 DEMOSTRACIÓN EN VIVO (1 minuto 45 segundos)', 1).runs[0].font.color.rgb = DUOC_BLUE
p = doc.add_paragraph()
run = p.add_run('📺 EJECUTAR EL PROGRAMA\n\n')
run.font.size = Pt(10)
run.bold = True
run.font.color.rgb = RGBColor(255, 165, 0)

p.add_run('Ahora les muestro el sistema en funcionamiento.\n\n').font.size = Pt(11)

# Opción 1
doc.add_heading('🔹 OPCIÓN 1: VENDER ENTRADA', 2).runs[0].font.color.rgb = DUOC_DARK_BLUE
demo_text = '''Selecciono la opción 1 para vender una entrada.
El sistema muestra la disponibilidad por sección.
Elijo ubicación VIP.
Veo los asientos disponibles del 1 al 30.
Selecciono el asiento 5.

Ingreso los datos del cliente:'''
p = doc.add_paragraph(demo_text)
p.runs[0].font.size = Pt(11)

datos_items = ['Nombre: Juan Pérez', 'RUT: 12345678', 'Edad: 25 años', '¿Es mujer? No', '¿Es estudiante? Sí']
for item in datos_items:
    p = doc.add_paragraph(item, style='List Bullet')
    p.runs[0].font.size = Pt(11)

p = doc.add_paragraph()
run = p.add_run('El sistema aplica automáticamente un descuento del 25% por estudiante.\n')
run.font.size = Pt(11)
run = p.add_run('Precio base 60 mil, descuento 15 mil, total a pagar 45 mil pesos.\n')
run.font.size = Pt(11)
run.bold = True
run = p.add_run('Venta exitosa.')
run.font.size = Pt(11)

# Opciones 4, 6, 7
doc.add_heading('🔹 OPCIÓN 4: IMPRIMIR BOLETA', 2).runs[0].font.color.rgb = DUOC_DARK_BLUE
p = doc.add_paragraph('Ahora imprimo la boleta del asiento 5.\nVemos toda la información: cliente, ubicación, fila, columna, y el desglose de precios.')
p.runs[0].font.size = Pt(11)

doc.add_heading('🔹 OPCIÓN 6: GENERAR REPORTE', 2).runs[0].font.color.rgb = DUOC_DARK_BLUE
p = doc.add_paragraph('Genero el reporte del sistema.\nMuestra que tengo 1 venta, 0.6% de ocupación, ingresos de 45 mil pesos, y el desglose por ubicación y tipo de cliente.')
p.runs[0].font.size = Pt(11)

doc.add_heading('🔹 OPCIÓN 7: EJECUTAR PRUEBAS', 2).runs[0].font.color.rgb = DUOC_DARK_BLUE
p = doc.add_paragraph('Ejecuto las pruebas del sistema.\nLas 4 pruebas pasan correctamente: integridad de datos OK, consistencia OK, validación de precios OK, y rendimiento en pocos milisegundos.')
p.runs[0].font.size = Pt(11)

# CARACTERÍSTICAS
doc.add_heading('✨ CARACTERÍSTICAS ÚNICAS (30 segundos)', 1).runs[0].font.color.rgb = DUOC_BLUE
p = doc.add_paragraph()
run = p.add_run('📺 MOSTRAR MÉTODO ejecutarPruebas() - LÍNEAS 511-559\n\n')
run.font.size = Pt(10)
run.bold = True
run.font.color.rgb = RGBColor(255, 165, 0)

p.add_run('Las características destacadas de este sistema son:\n\n').font.size = Pt(11)

caract_items = [
    'Suite completa de 4 pruebas automatizadas que validan integridad, consistencia, precios y rendimiento.',
    'Sistema inteligente de descuentos que detecta automáticamente la edad del cliente.',
    'Funcionalidad de modificación de asientos con recálculo automático de precios cuando cambia la ubicación.',
    'Validaciones robustas que previenen errores y duplicados.',
    'Visualización gráfica del estado de todos los asientos por sección.'
]

for item in caract_items:
    p = doc.add_paragraph(item, style='List Bullet')
    p.runs[0].font.size = Pt(11)

# CIERRE
doc.add_heading('🎯 CIERRE (15 segundos)', 1).runs[0].font.color.rgb = DUOC_BLUE
p = doc.add_paragraph()
run = p.add_run('📺 SELECCIONAR OPCIÓN 8: SALIR\n\n')
run.font.size = Pt(10)
run.bold = True
run.font.color.rgb = RGBColor(255, 165, 0)

p.add_run('Al salir, el sistema muestra el resumen: total vendido y cantidad de ventas.\n\n').font.size = Pt(11)
run = p.add_run('Esto concluye la presentación del Sistema Teatro Moro.\nMuchas gracias por su atención.')
run.font.size = Pt(11)
run.bold = True

doc.add_page_break()

# TABLA DE TIEMPO
doc.add_heading('📊 DISTRIBUCIÓN DEL TIEMPO', 1).runs[0].font.color.rgb = DUOC_BLUE
table = doc.add_table(rows=8, cols=3)
table.style = 'Light Grid Accent 1'

# Encabezados
headers = table.rows[0].cells
headers[0].text = 'Sección'
headers[1].text = 'Duración'
headers[2].text = 'Tiempo Acumulado'

# Datos
data = [
    ('Introducción', '0:30', '0:30'),
    ('Arquitectura del Sistema', '0:45', '1:15'),
    ('Implementación y Estrategias', '1:00', '2:15'),
    ('Dificultades y Soluciones', '0:45', '3:00'),
    ('Demostración en Vivo', '1:45', '4:45'),
    ('Características Únicas', '0:30', '5:15*'),
    ('Cierre', '0:15', '5:00')
]

for i, (sec, dur, acum) in enumerate(data, start=1):
    cells = table.rows[i].cells
    cells[0].text = sec
    cells[1].text = dur
    cells[2].text = acum

# TIPS
doc.add_heading('💡 TIPS PARA LA GRABACIÓN', 1).runs[0].font.color.rgb = DUOC_BLUE
tips = [
    'Practica 2-3 veces antes de grabar',
    'Habla natural, no como robot',
    'Haz pausas breves entre secciones',
    'Sincroniza lo que dices con lo que muestras en pantalla',
    'Si te equivocas, respira y continúa',
    'Mira a la cámara cuando hablas (no leas directamente)',
    'Sonríe al inicio y al final'
]

for tip in tips:
    p = doc.add_paragraph(tip, style='List Number')
    p.runs[0].font.size = Pt(11)

# Guardar Word
word_path = 'GUION_VIDEO_ETF_S9_Fuad_Onate.docx'
doc.save(word_path)
print(f"✅ Documento Word creado: {word_path}")

# ==================== CREAR PRESENTACIÓN POWERPOINT ====================
print("\n2. Creando presentación PowerPoint...")
prs = Presentation()
prs.slide_width = PptInches(10)
prs.slide_height = PptInches(7.5)

# SLIDE 1: PORTADA
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide
left = PptInches(0.5)
top = PptInches(2)
width = PptInches(9)
height = PptInches(1.5)

# Título
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.text = "🎬 GUION VIDEO - SISTEMA TEATRO MORO"
tf.paragraphs[0].font.size = PptPt(44)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = PptRGBColor(0, 102, 204)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Subtítulo
txBox2 = slide.shapes.add_textbox(left, PptInches(3.5), width, PptInches(2))
tf2 = txBox2.text_frame
tf2.text = "Evaluación Final Transversal S9\nFuad Oñate - Programación Online\nFundamentos de Programación"
for paragraph in tf2.paragraphs:
    paragraph.font.size = PptPt(24)
    paragraph.alignment = PP_ALIGN.CENTER

# SLIDE 2: INTRODUCCIÓN
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "📍 INTRODUCCIÓN (30 seg)"
title.text_frame.paragraphs[0].font.color.rgb = PptRGBColor(0, 102, 204)

content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Hola, mi nombre es Fuad Oñate y les presento el Sistema de Gestión de Ventas de Entradas para el Teatro Moro."
p = tf.add_paragraph()
p.text = "Sistema desarrollado para EFT - Fundamentos de Programación, Semana 9"
p.level = 1

# SLIDE 3: ARQUITECTURA
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "🏗️ ARQUITECTURA DEL SISTEMA (45 seg)"
title.text_frame.paragraphs[0].font.color.rgb = PptRGBColor(0, 102, 204)

content = slide.placeholders[1]
tf = content.text_frame
tf.text = "10 arrays paralelos (150 asientos)"

for item in ["ArrayList para historial dinámico", "Variables estáticas globales", "5 constantes de precios", "5 secciones × 30 asientos"]:
    p = tf.add_paragraph()
    p.text = item
    p.level = 1

# SLIDE 4: IMPLEMENTACIÓN
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "⚙️ IMPLEMENTACIÓN (1 min)"
title.text_frame.paragraphs[0].font.color.rgb = PptRGBColor(0, 102, 204)

content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Estructuras de Control:"

for item in ["IF-ELSE → Validaciones y descuentos", "SWITCH → Menú y ubicaciones", "FOR → Recorrer asientos", "DO-WHILE → Menú activo", "TRY-CATCH → Manejo errores"]:
    p = tf.add_paragraph()
    p.text = item
    p.level = 1

# SLIDE 5: DIFICULTADES
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "🔧 DIFICULTADES Y SOLUCIONES (45 seg)"
title.text_frame.paragraphs[0].font.color.rgb = PptRGBColor(0, 102, 204)

content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Desafío: Recalcular precios al cambiar ubicación"

p = tf.add_paragraph()
p.text = "Solución: Mantener % descuento y aplicar a nuevo precio"
p.level = 1

p = tf.add_paragraph()
p.text = "Desafío: Validar RUT duplicado"
p.level = 0

p = tf.add_paragraph()
p.text = "Solución: Recorrer asientos vendidos y comparar"
p.level = 1

# SLIDE 6: DEMO
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "💻 DEMOSTRACIÓN (1 min 45 seg)"
title.text_frame.paragraphs[0].font.color.rgb = PptRGBColor(0, 102, 204)

content = slide.placeholders[1]
tf = content.text_frame
tf.text = "1️⃣ Vender Entrada (Juan Pérez, estudiante)"

for item in ["4️⃣ Imprimir Boleta", "6️⃣ Generar Reporte", "7️⃣ Ejecutar Pruebas (4 tests OK)"]:
    p = tf.add_paragraph()
    p.text = item
    p.level = 0

# SLIDE 7: CARACTERÍSTICAS
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "✨ CARACTERÍSTICAS ÚNICAS (30 seg)"
title.text_frame.paragraphs[0].font.color.rgb = PptRGBColor(0, 102, 204)

content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Suite de 4 pruebas automatizadas"

for item in ["Descuentos automáticos por edad", "Recálculo inteligente de precios", "Validaciones robustas", "Visualización gráfica de asientos"]:
    p = tf.add_paragraph()
    p.text = item
    p.level = 1

# SLIDE 8: DISTRIBUCIÓN TIEMPO
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "📊 DISTRIBUCIÓN DEL TIEMPO"
title.text_frame.paragraphs[0].font.color.rgb = PptRGBColor(0, 102, 204)

# Tabla
rows, cols = 8, 3
left = PptInches(1)
top = PptInches(2)
width = PptInches(8)
height = PptInches(4)

table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Encabezados
table.cell(0, 0).text = "Sección"
table.cell(0, 1).text = "Duración"
table.cell(0, 2).text = "Acumulado"

# Datos
time_data = [
    ("Introducción", "0:30", "0:30"),
    ("Arquitectura", "0:45", "1:15"),
    ("Implementación", "1:00", "2:15"),
    ("Dificultades", "0:45", "3:00"),
    ("Demostración", "1:45", "4:45"),
    ("Características", "0:30", "5:15*"),
    ("Cierre", "0:15", "5:00")
]

for i, (sec, dur, acum) in enumerate(time_data, start=1):
    table.cell(i, 0).text = sec
    table.cell(i, 1).text = dur
    table.cell(i, 2).text = acum

# SLIDE 9: CIERRE
slide = prs.slides.add_slide(prs.slide_layouts[6])
txBox = slide.shapes.add_textbox(PptInches(1), PptInches(2.5), PptInches(8), PptInches(2))
tf = txBox.text_frame
tf.text = "🎯 ¡Éxito con tu grabación!"
tf.paragraphs[0].font.size = PptPt(54)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = PptRGBColor(0, 102, 204)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = "Fuad Oñate - Sistema Teatro Moro"
p.font.size = PptPt(28)
p.alignment = PP_ALIGN.CENTER

# Guardar PowerPoint
ppt_path = 'GUION_VIDEO_ETF_S9_Fuad_Onate.pptx'
prs.save(ppt_path)
print(f"✅ Presentación PowerPoint creada: {ppt_path}")

print("\n" + "="*60)
print("🎉 ¡ARCHIVOS CREADOS EXITOSAMENTE!")
print("="*60)
print(f"\n📄 Word:  {word_path}")
print(f"📊 PowerPoint: {ppt_path}")
print("\nUbicación: s9/")
