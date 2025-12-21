# -*- coding: utf-8 -*-
import sys
sys.stdout.reconfigure(encoding='utf-8')

from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PptRGBColor

# Colores Duoc
DUOC_BLUE = RGBColor(0, 102, 204)
DUOC_DARK_BLUE = RGBColor(0, 68, 153)

print('Creando documentos profesionales...')

# ==================== DOCUMENTO WORD COMPLETO ====================
print('\\n1. Creando documento Word completo...')
doc = Document()

# Configurar márgenes
for section in doc.sections:
    section.top_margin = Inches(1)
    section.bottom_margin = Inches(1)
    section.left_margin = Inches(1)
    section.right_margin = Inches(1)

# PORTADA
h = doc.add_heading('GUION PARA VIDEO - EVALUACION FINAL TRANSVERSAL', 0)
h.alignment = WD_ALIGN_PARAGRAPH.CENTER
h.runs[0].font.color.rgb = DUOC_BLUE

h2 = doc.add_heading('Sistema Teatro Moro - Semana 9', 2)
h2.alignment = WD_ALIGN_PARAGRAPH.CENTER
h2.runs[0].font.color.rgb = DUOC_DARK_BLUE

# Info del encabezado
info = doc.add_paragraph()
info.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = info.add_run('Fundamentos de Programacion (PRY2201)\\n')
run.font.size = Pt(14)
run.bold = True
run = info.add_run('Duracion: 5 minutos maximo\\n')
run.font.size = Pt(12)
run = info.add_run('\\n')
run = info.add_run('Estudiante: Fuad Onate\\n')
run.font.size = Pt(14)
run.bold = True
run.font.color.rgb = DUOC_BLUE
run = info.add_run('Programacion Online - Duoc UC')
run.font.size = Pt(12)

doc.add_page_break()

# INTRODUCCION
doc.add_heading('INTRODUCCION (30 segundos)', 1).runs[0].font.color.rgb = DUOC_BLUE
p = doc.add_paragraph()
p.add_run('Hola, mi nombre es Fuad Onate y les presento el Sistema de Gestion de Ventas de Entradas para el Teatro Moro, desarrollado como Evaluacion Final Transversal de la asignatura Fundamentos de Programacion, semana 9.\\n\\n').font.size = Pt(11)
p.add_run('Este sistema permite vender entradas, gestionar asientos y aplicar descuentos automaticos segun el tipo de cliente.').font.size = Pt(11)

# ARQUITECTURA
doc.add_heading('ARQUITECTURA DEL SISTEMA (45 segundos)', 1).runs[0].font.color.rgb = DUOC_BLUE
p = doc.add_paragraph()
run = p.add_run('[MOSTRAR CODIGO EN NETBEANS - LINEAS 14-39]\\n\\n')
run.font.size = Pt(10)
run.bold = True
run.font.color.rgb = RGBColor(255, 140, 0)

p.add_run('El sistema esta construido con:\\n\\n').font.size = Pt(11)

for item in ['10 arrays paralelos para 150 asientos', 'ArrayList para historial dinamico', 'Variables estaticas globales', '5 constantes de precios', '5 secciones de 30 asientos']:
    doc.add_paragraph(item, style='List Bullet').runs[0].font.size = Pt(11)

# IMPLEMENTACION
doc.add_heading('IMPLEMENTACION Y ESTRATEGIAS (1 minuto)', 1).runs[0].font.color.rgb = DUOC_BLUE
p = doc.add_paragraph()
run = p.add_run('[MOSTRAR METODO venderEntrada()]\\n\\n')
run.font.size = Pt(10)
run.bold = True
run.font.color.rgb = RGBColor(255, 140, 0)

p.add_run('Estructuras de control implementadas:\\n\\n').font.size = Pt(11)

for item in ['IF-ELSE para validaciones y descuentos', 'SWITCH para menu y ubicaciones', 'FOR para recorrer asientos', 'DO-WHILE para menu activo', 'TRY-CATCH para excepciones']:
    doc.add_paragraph(item, style='List Bullet').runs[0].font.size = Pt(11)

# DIFICULTADES
doc.add_heading('DIFICULTADES Y SOLUCIONES (45 segundos)', 1).runs[0].font.color.rgb = DUOC_BLUE
p = doc.add_paragraph()
run = p.add_run('[MOSTRAR METODO modificarAsiento() - LINEAS 304-317]\\n\\n')
run.font.size = Pt(10)
run.bold = True
run.font.color.rgb = RGBColor(255, 140, 0)

p.add_run('Principales desafios:\\n\\n1. Recalculo de precios al cambiar ubicacion\\nSolucion: Mantener porcentaje de descuento\\n\\n2. Validacion de RUT duplicado\\nSolucion: Recorrer asientos vendidos').font.size = Pt(11)

doc.add_page_break()

# DEMOSTRACION
doc.add_heading('DEMOSTRACION EN VIVO (1 minuto 45 segundos)', 1).runs[0].font.color.rgb = DUOC_BLUE
p = doc.add_paragraph()
run = p.add_run('[EJECUTAR EL PROGRAMA]\\n\\n')
run.font.size = Pt(10)
run.bold = True
run.font.color.rgb = RGBColor(255, 140, 0)

doc.add_heading('Opcion 1: Vender Entrada', 2).runs[0].font.color.rgb = DUOC_DARK_BLUE
p = doc.add_paragraph('- Ubicacion: VIP\\n- Asiento: 5\\n- Cliente: Juan Perez (RUT: 12345678)\\n- Edad: 25 anos\\n- Estudiante: Si\\n- Descuento: 25% = $15,000\\n- Total: $45,000')
p.runs[0].font.size = Pt(11)

doc.add_heading('Opcion 4: Imprimir Boleta', 2).runs[0].font.color.rgb = DUOC_DARK_BLUE
doc.add_paragraph('Muestra informacion completa del cliente y desglose de precios').runs[0].font.size = Pt(11)

doc.add_heading('Opcion 6: Generar Reporte', 2).runs[0].font.color.rgb = DUOC_DARK_BLUE
doc.add_paragraph('Estadisticas: 1 venta, 0.6% ocupacion, $45,000 ingresos').runs[0].font.size = Pt(11)

doc.add_heading('Opcion 7: Ejecutar Pruebas', 2).runs[0].font.color.rgb = DUOC_DARK_BLUE
doc.add_paragraph('4 pruebas OK: Integridad, Consistencia, Precios, Rendimiento').runs[0].font.size = Pt(11)

# CARACTERISTICAS
doc.add_heading('CARACTERISTICAS UNICAS (30 segundos)', 1).runs[0].font.color.rgb = DUOC_BLUE
p = doc.add_paragraph()
run = p.add_run('[MOSTRAR METODO ejecutarPruebas() - LINEAS 511-559]\\n\\n')
run.font.size = Pt(10)
run.bold = True
run.font.color.rgb = RGBColor(255, 140, 0)

for item in ['Suite de 4 pruebas automatizadas', 'Descuentos automaticos por edad', 'Recalculo inteligente de precios', 'Validaciones robustas', 'Visualizacion grafica']:
    doc.add_paragraph(item, style='List Bullet').runs[0].font.size = Pt(11)

# CIERRE
doc.add_heading('CIERRE (15 segundos)', 1).runs[0].font.color.rgb = DUOC_BLUE
p = doc.add_paragraph()
run = p.add_run('[SELECCIONAR OPCION 8: SALIR]\\n\\n')
run.font.size = Pt(10)
run.bold = True
run.font.color.rgb = RGBColor(255, 140, 0)
p.add_run('Al salir muestra resumen final.\\n\\n')
run = p.add_run('Esto concluye la presentacion del Sistema Teatro Moro.\\nMuchas gracias por su atencion.')
run.font.size = Pt(12)
run.bold = True

doc.add_page_break()

# TABLA DE TIEMPO
doc.add_heading('DISTRIBUCION DEL TIEMPO', 1).runs[0].font.color.rgb = DUOC_BLUE
table = doc.add_table(rows=8, cols=3)
table.style = 'Light Grid Accent 1'

headers = table.rows[0].cells
headers[0].text = 'Seccion'
headers[1].text = 'Duracion'
headers[2].text = 'Acumulado'

data = [
    ('Introduccion', '0:30', '0:30'),
    ('Arquitectura', '0:45', '1:15'),
    ('Implementacion', '1:00', '2:15'),
    ('Dificultades', '0:45', '3:00'),
    ('Demostracion', '1:45', '4:45'),
    ('Caracteristicas', '0:30', '5:15*'),
    ('Cierre', '0:15', '5:00')
]

for i, (sec, dur, acum) in enumerate(data, start=1):
    cells = table.rows[i].cells
    cells[0].text = sec
    cells[1].text = dur
    cells[2].text = acum

# TIPS
doc.add_heading('TIPS PARA LA GRABACION', 1).runs[0].font.color.rgb = DUOC_BLUE
tips = ['Practica 2-3 veces', 'Habla natural', 'Pausas breves', 'Sincroniza audio/video', 'Respira si te equivocas', 'Mira a la camara', 'Sonrie']
for tip in tips:
    doc.add_paragraph(tip, style='List Number').runs[0].font.size = Pt(11)

# FOOTER
p = doc.add_paragraph()
p.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = p.add_run('\\n\\nExito con tu grabacion!\\nFuad Onate - Programacion Online - Duoc UC')
run.font.size = Pt(14)
run.bold = True
run.font.color.rgb = DUOC_BLUE

word_file = 'GUION_VIDEO_COMPLETO_Fuad_Onate.docx'
doc.save(word_file)
print(f'OK - Word guardado: {word_file}')

# ==================== POWERPOINT COMPLETO ====================
print('\\n2. Creando PowerPoint completo...')
prs = Presentation()
prs.slide_width = PptInches(10)
prs.slide_height = PptInches(7.5)

# SLIDE 1: PORTADA
slide = prs.slides.add_slide(prs.slide_layouts[6])
txBox = slide.shapes.add_textbox(PptInches(0.5), PptInches(2), PptInches(9), PptInches(1.5))
tf = txBox.text_frame
tf.text = 'GUION VIDEO - SISTEMA TEATRO MORO'
tf.paragraphs[0].font.size = PptPt(44)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = PptRGBColor(0, 102, 204)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

txBox2 = slide.shapes.add_textbox(PptInches(0.5), PptInches(3.5), PptInches(9), PptInches(2))
tf2 = txBox2.text_frame
tf2.text = 'EFT Semana 9\\nFuad Onate - Programacion Online\\nDuoc UC'
for p in tf2.paragraphs:
    p.font.size = PptPt(24)
    p.alignment = PP_ALIGN.CENTER

# SLIDE 2: INTRODUCCION
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'INTRODUCCION (30 seg)'
slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = PptRGBColor(0, 102, 204)
tf = slide.placeholders[1].text_frame
tf.text = 'Sistema de Gestion de Ventas de Entradas'
p = tf.add_paragraph()
p.text = 'Teatro Moro - EFT Semana 9'
p.level = 1

# SLIDE 3: ARQUITECTURA
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'ARQUITECTURA (45 seg)'
slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = PptRGBColor(0, 102, 204)
tf = slide.placeholders[1].text_frame
tf.text = '10 arrays paralelos (150 asientos)'
for item in ['ArrayList dinamico', 'Variables estaticas', '5 constantes precios', '5 secciones x 30']:
    p = tf.add_paragraph()
    p.text = item
    p.level = 1

# SLIDE 4: IMPLEMENTACION
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'IMPLEMENTACION (1 min)'
slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = PptRGBColor(0, 102, 204)
tf = slide.placeholders[1].text_frame
tf.text = 'Estructuras de Control:'
for item in ['IF-ELSE: Validaciones', 'SWITCH: Menu', 'FOR: Recorrer', 'DO-WHILE: Loop', 'TRY-CATCH: Errores']:
    p = tf.add_paragraph()
    p.text = item
    p.level = 1

# SLIDE 5: DEMO
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'DEMOSTRACION (1 min 45 seg)'
slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = PptRGBColor(0, 102, 204)
tf = slide.placeholders[1].text_frame
tf.text = 'Vender Entrada (Juan Perez, estudiante)'
for item in ['Imprimir Boleta', 'Generar Reporte', 'Ejecutar Pruebas (4 OK)']:
    p = tf.add_paragraph()
    p.text = item
    p.level = 0

# SLIDE 6: CARACTERISTICAS
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'CARACTERISTICAS UNICAS (30 seg)'
slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = PptRGBColor(0, 102, 204)
tf = slide.placeholders[1].text_frame
tf.text = '4 pruebas automatizadas'
for item in ['Descuentos automaticos', 'Recalculo inteligente', 'Validaciones robustas']:
    p = tf.add_paragraph()
    p.text = item
    p.level = 1

# SLIDE 7: TIEMPO
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = 'DISTRIBUCION DEL TIEMPO'
slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = PptRGBColor(0, 102, 204)
table = slide.shapes.add_table(8, 3, PptInches(1), PptInches(2), PptInches(8), PptInches(4)).table
table.cell(0, 0).text = 'Seccion'
table.cell(0, 1).text = 'Duracion'
table.cell(0, 2).text = 'Acumulado'
for i, (sec, dur, acum) in enumerate(data, start=1):
    table.cell(i, 0).text = sec
    table.cell(i, 1).text = dur
    table.cell(i, 2).text = acum

# SLIDE 8: CIERRE
slide = prs.slides.add_slide(prs.slide_layouts[6])
txBox = slide.shapes.add_textbox(PptInches(1), PptInches(2.5), PptInches(8), PptInches(2))
tf = txBox.text_frame
tf.text = 'Exito con tu grabacion!'
tf.paragraphs[0].font.size = PptPt(54)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = PptRGBColor(0, 102, 204)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = 'Fuad Onate - Duoc UC'
p.font.size = PptPt(28)
p.alignment = PP_ALIGN.CENTER

ppt_file = 'GUION_VIDEO_COMPLETO_Fuad_Onate.pptx'
prs.save(ppt_file)
print(f'OK - PowerPoint guardado: {ppt_file}')

print('\\n' + '='*50)
print('ARCHIVOS CREADOS EXITOSAMENTE!')
print('='*50)
print(f'Word: {word_file}')
print(f'PowerPoint: {ppt_file}')
