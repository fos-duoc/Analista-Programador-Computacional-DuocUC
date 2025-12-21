# -*- coding: utf-8 -*-
import sys
sys.stdout.reconfigure(encoding='utf-8')

from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PptRGBColor

# Colores Duoc
DUOC_BLUE = RGBColor(0, 102, 204)
DUOC_DARK_BLUE = RGBColor(0, 68, 153)
ORANGE = RGBColor(255, 140, 0)

print('Creando documentos profesionales con iconos...')

# ==================== DOCUMENTO WORD COMPLETO ====================
print('\n📄 Creando documento Word completo...')
doc = Document()

# Configurar márgenes
for section in doc.sections:
    section.top_margin = Inches(1)
    section.bottom_margin = Inches(1)
    section.left_margin = Inches(1)
    section.right_margin = Inches(1)

# PORTADA
h = doc.add_heading('🎬 GUION PARA VIDEO - EVALUACIÓN FINAL TRANSVERSAL', 0)
h.alignment = WD_ALIGN_PARAGRAPH.CENTER
h.runs[0].font.color.rgb = DUOC_BLUE

h2 = doc.add_heading('💻 Sistema Teatro Moro - Semana 9', 2)
h2.alignment = WD_ALIGN_PARAGRAPH.CENTER
h2.runs[0].font.color.rgb = DUOC_DARK_BLUE

# Info del encabezado
info = doc.add_paragraph()
info.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = info.add_run('☕ Fundamentos de Programación (PRY2201)\n')
run.font.size = Pt(14)
run.bold = True
run = info.add_run('⏱️ Duración: 5 minutos máximo\n')
run.font.size = Pt(12)
run = info.add_run('\n')
run = info.add_run('👨‍💻 Estudiante: Fuad Oñate\n')
run.font.size = Pt(14)
run.bold = True
run.font.color.rgb = DUOC_BLUE
run = info.add_run('🎓 Programación Online - Duoc UC')
run.font.size = Pt(12)

doc.add_page_break()

# INTRODUCCIÓN
doc.add_heading('📍 INTRODUCCIÓN (30 segundos)', 1).runs[0].font.color.rgb = DUOC_BLUE
p = doc.add_paragraph()
p.add_run('Hola, mi nombre es Fuad Oñate y les presento el Sistema de Gestión de Ventas de Entradas para el Teatro Moro, desarrollado como Evaluación Final Transversal de la asignatura Fundamentos de Programación, semana 9.\n\n').font.size = Pt(11)
p.add_run('Este sistema permite vender entradas, gestionar asientos y aplicar descuentos automáticos según el tipo de cliente.').font.size = Pt(11)

# ARQUITECTURA
doc.add_heading('🏗️ ARQUITECTURA DEL SISTEMA (45 segundos)', 1).runs[0].font.color.rgb = DUOC_BLUE
p = doc.add_paragraph()
run = p.add_run('📺 [MOSTRAR CÓDIGO EN NETBEANS - LÍNEAS 14-39]\n\n')
run.font.size = Pt(10)
run.bold = True
run.font.color.rgb = ORANGE

p.add_run('El sistema está construido con:\n\n').font.size = Pt(11)

items_arq = [
    '📊 10 arrays paralelos para 150 asientos',
    '📋 ArrayList para historial dinámico',
    '🔢 Variables estáticas globales',
    '💰 5 constantes de precios',
    '🎭 5 secciones de 30 asientos'
]

for item in items_arq:
    doc.add_paragraph(item, style='List Bullet').runs[0].font.size = Pt(11)

# IMPLEMENTACIÓN
doc.add_heading('⚙️ IMPLEMENTACIÓN Y ESTRATEGIAS (1 minuto)', 1).runs[0].font.color.rgb = DUOC_BLUE
p = doc.add_paragraph()
run = p.add_run('📺 [MOSTRAR MÉTODO venderEntrada()]\n\n')
run.font.size = Pt(10)
run.bold = True
run.font.color.rgb = ORANGE

p.add_run('Estructuras de control implementadas:\n\n').font.size = Pt(11)

items_impl = [
    '🔀 IF-ELSE para validaciones y descuentos',
    '🔄 SWITCH para menú y ubicaciones',
    '🔁 FOR para recorrer asientos',
    '♾️ DO-WHILE para menú activo',
    '⚠️ TRY-CATCH para excepciones'
]

for item in items_impl:
    doc.add_paragraph(item, style='List Bullet').runs[0].font.size = Pt(11)

# DIFICULTADES
doc.add_heading('🔧 DIFICULTADES Y SOLUCIONES (45 segundos)', 1).runs[0].font.color.rgb = DUOC_BLUE
p = doc.add_paragraph()
run = p.add_run('📺 [MOSTRAR MÉTODO modificarAsiento() - LÍNEAS 304-317]\n\n')
run.font.size = Pt(10)
run.bold = True
run.font.color.rgb = ORANGE

p.add_run('Principales desafíos:\n\n').font.size = Pt(11)
p.add_run('❌ Problema 1: Recálculo de precios al cambiar ubicación\n').font.size = Pt(11)
p.add_run('✅ Solución: Mantener porcentaje de descuento\n\n').font.size = Pt(11)
p.add_run('❌ Problema 2: Validación de RUT duplicado\n').font.size = Pt(11)
p.add_run('✅ Solución: Recorrer asientos vendidos').font.size = Pt(11)

doc.add_page_break()

# DEMOSTRACIÓN
doc.add_heading('💻 DEMOSTRACIÓN EN VIVO (1 minuto 45 segundos)', 1).runs[0].font.color.rgb = DUOC_BLUE
p = doc.add_paragraph()
run = p.add_run('📺 [EJECUTAR EL PROGRAMA]\n\n')
run.font.size = Pt(10)
run.bold = True
run.font.color.rgb = ORANGE

doc.add_heading('1️⃣ Opción 1: Vender Entrada', 2).runs[0].font.color.rgb = DUOC_DARK_BLUE
p = doc.add_paragraph('🎯 Ubicación: VIP\n🪑 Asiento: 5\n👤 Cliente: Juan Pérez (RUT: 12345678)\n🎂 Edad: 25 años\n📚 Estudiante: Sí\n💸 Descuento: 25% = $15,000\n💵 Total: $45,000')
p.runs[0].font.size = Pt(11)

doc.add_heading('4️⃣ Opción 4: Imprimir Boleta', 2).runs[0].font.color.rgb = DUOC_DARK_BLUE
doc.add_paragraph('🧾 Muestra información completa del cliente y desglose de precios').runs[0].font.size = Pt(11)

doc.add_heading('6️⃣ Opción 6: Generar Reporte', 2).runs[0].font.color.rgb = DUOC_DARK_BLUE
doc.add_paragraph('📈 Estadísticas: 1 venta, 0.6% ocupación, $45,000 ingresos').runs[0].font.size = Pt(11)

doc.add_heading('7️⃣ Opción 7: Ejecutar Pruebas', 2).runs[0].font.color.rgb = DUOC_DARK_BLUE
doc.add_paragraph('✅ 4 pruebas OK: Integridad, Consistencia, Precios, Rendimiento').runs[0].font.size = Pt(11)

# CARACTERÍSTICAS
doc.add_heading('✨ CARACTERÍSTICAS ÚNICAS (30 segundos)', 1).runs[0].font.color.rgb = DUOC_BLUE
p = doc.add_paragraph()
run = p.add_run('📺 [MOSTRAR MÉTODO ejecutarPruebas() - LÍNEAS 511-559]\n\n')
run.font.size = Pt(10)
run.bold = True
run.font.color.rgb = ORANGE

items_caract = [
    '🧪 Suite de 4 pruebas automatizadas',
    '🎁 Descuentos automáticos por edad',
    '🔄 Recálculo inteligente de precios',
    '🛡️ Validaciones robustas',
    '📊 Visualización gráfica'
]

for item in items_caract:
    doc.add_paragraph(item, style='List Bullet').runs[0].font.size = Pt(11)

# CIERRE
doc.add_heading('🎯 CIERRE (15 segundos)', 1).runs[0].font.color.rgb = DUOC_BLUE
p = doc.add_paragraph()
run = p.add_run('📺 [SELECCIONAR OPCIÓN 8: SALIR]\n\n')
run.font.size = Pt(10)
run.bold = True
run.font.color.rgb = ORANGE
p.add_run('Al salir muestra resumen final.\n\n').font.size = Pt(11)
run = p.add_run('Esto concluye la presentación del Sistema Teatro Moro.\nMuchas gracias por su atención.')
run.font.size = Pt(12)
run.bold = True

doc.add_page_break()

# TABLA DE TIEMPO
doc.add_heading('📊 DISTRIBUCIÓN DEL TIEMPO', 1).runs[0].font.color.rgb = DUOC_BLUE
table = doc.add_table(rows=8, cols=3)
table.style = 'Light Grid Accent 1'

headers = table.rows[0].cells
headers[0].text = 'Sección'
headers[1].text = 'Duración'
headers[2].text = 'Acumulado'

data = [
    ('📍 Introducción', '0:30', '0:30'),
    ('🏗️ Arquitectura', '0:45', '1:15'),
    ('⚙️ Implementación', '1:00', '2:15'),
    ('🔧 Dificultades', '0:45', '3:00'),
    ('💻 Demostración', '1:45', '4:45'),
    ('✨ Características', '0:30', '5:15*'),
    ('🎯 Cierre', '0:15', '5:00')
]

for i, (sec, dur, acum) in enumerate(data, start=1):
    cells = table.rows[i].cells
    cells[0].text = sec
    cells[1].text = dur
    cells[2].text = acum

# TIPS
doc.add_heading('💡 TIPS PARA LA GRABACIÓN', 1).runs[0].font.color.rgb = DUOC_BLUE
tips = [
    '🔄 Practica 2-3 veces',
    '🗣️ Habla natural',
    '⏸️ Pausas breves',
    '🎬 Sincroniza audio/video',
    '😌 Respira si te equivocas',
    '📷 Mira a la cámara',
    '😊 Sonríe'
]
for tip in tips:
    doc.add_paragraph(tip, style='List Number').runs[0].font.size = Pt(11)

# ACCIONES EN PANTALLA
doc.add_heading('🎬 ACCIONES EN PANTALLA', 1).runs[0].font.color.rgb = DUOC_BLUE
table2 = doc.add_table(rows=6, cols=2)
table2.style = 'Light Grid Accent 1'

headers2 = table2.rows[0].cells
headers2[0].text = 'Sección'
headers2[1].text = 'Qué Mostrar'

acciones = [
    ('🏗️ Arquitectura', 'Líneas 14-39 (arrays y variables)'),
    ('⚙️ Implementación', 'Método venderEntrada() - líneas 80-258'),
    ('🔧 Dificultades', 'Método modificarAsiento() - líneas 304-317'),
    ('💻 Demostración', 'Ejecutar programa en consola'),
    ('✨ Características', 'Método ejecutarPruebas() - líneas 511-559')
]

for i, (sec, mostrar) in enumerate(acciones, start=1):
    cells = table2.rows[i].cells
    cells[0].text = sec
    cells[1].text = mostrar

# FOOTER
p = doc.add_paragraph()
p.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = p.add_run('\n\n🎉 ¡Éxito con tu grabación!\n👨‍💻 Fuad Oñate - Programación Online\n🎓 Duoc UC')
run.font.size = Pt(14)
run.bold = True
run.font.color.rgb = DUOC_BLUE

word_file = 'GUION_VIDEO_FINAL_Fuad_Onate.docx'
doc.save(word_file)
print(f'✅ Word guardado: {word_file}')

# ==================== POWERPOINT COMPLETO ====================
print('\n📊 Creando PowerPoint completo...')
prs = Presentation()
prs.slide_width = PptInches(10)
prs.slide_height = PptInches(7.5)

# SLIDE 1: PORTADA
slide = prs.slides.add_slide(prs.slide_layouts[6])
txBox = slide.shapes.add_textbox(PptInches(0.5), PptInches(2), PptInches(9), PptInches(1.5))
tf = txBox.text_frame
tf.text = '🎬 GUION VIDEO\nSISTEMA TEATRO MORO'
tf.paragraphs[0].font.size = PptPt(44)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = PptRGBColor(0, 102, 204)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

txBox2 = slide.shapes.add_textbox(PptInches(0.5), PptInches(4), PptInches(9), PptInches(2))
tf2 = txBox2.text_frame
tf2.text = '📚 EFT Semana 9\n👨‍💻 Fuad Oñate\n🎓 Programación Online - Duoc UC'
for p in tf2.paragraphs:
    p.font.size = PptPt(24)
    p.alignment = PP_ALIGN.CENTER

# SLIDE 2: INTRODUCCIÓN
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = '📍 INTRODUCCIÓN (30 seg)'
slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = PptRGBColor(0, 102, 204)
tf = slide.placeholders[1].text_frame
tf.text = '💻 Sistema de Gestión de Ventas'
p = tf.add_paragraph()
p.text = '🎭 Teatro Moro - EFT Semana 9'
p.level = 1

# SLIDE 3: ARQUITECTURA
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = '🏗️ ARQUITECTURA (45 seg)'
slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = PptRGBColor(0, 102, 204)
tf = slide.placeholders[1].text_frame
tf.text = '📊 10 arrays paralelos (150 asientos)'
for item in ['📋 ArrayList dinámico', '🔢 Variables estáticas', '💰 5 constantes precios', '🎭 5 secciones × 30']:
    p = tf.add_paragraph()
    p.text = item
    p.level = 1

# SLIDE 4: IMPLEMENTACIÓN
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = '⚙️ IMPLEMENTACIÓN (1 min)'
slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = PptRGBColor(0, 102, 204)
tf = slide.placeholders[1].text_frame
tf.text = '🔧 Estructuras de Control:'
for item in ['🔀 IF-ELSE: Validaciones', '🔄 SWITCH: Menú', '🔁 FOR: Recorrer', '♾️ DO-WHILE: Loop', '⚠️ TRY-CATCH: Errores']:
    p = tf.add_paragraph()
    p.text = item
    p.level = 1

# SLIDE 5: DIFICULTADES
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = '🔧 DIFICULTADES (45 seg)'
slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = PptRGBColor(0, 102, 204)
tf = slide.placeholders[1].text_frame
tf.text = '❌ Recálculo de precios'
p = tf.add_paragraph()
p.text = '✅ Mantener % descuento'
p.level = 1
p = tf.add_paragraph()
p.text = '❌ RUT duplicado'
p.level = 0
p = tf.add_paragraph()
p.text = '✅ Validación en array'
p.level = 1

# SLIDE 6: DEMO
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = '💻 DEMOSTRACIÓN (1:45)'
slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = PptRGBColor(0, 102, 204)
tf = slide.placeholders[1].text_frame
tf.text = '1️⃣ Vender (Juan Pérez, estudiante)'
for item in ['4️⃣ Imprimir Boleta', '6️⃣ Generar Reporte', '7️⃣ Pruebas (4 ✅ OK)']:
    p = tf.add_paragraph()
    p.text = item
    p.level = 0

# SLIDE 7: CARACTERÍSTICAS
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = '✨ CARACTERÍSTICAS (30 seg)'
slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = PptRGBColor(0, 102, 204)
tf = slide.placeholders[1].text_frame
tf.text = '🧪 4 pruebas automatizadas'
for item in ['🎁 Descuentos automáticos', '🔄 Recálculo inteligente', '🛡️ Validaciones robustas', '📊 Visualización gráfica']:
    p = tf.add_paragraph()
    p.text = item
    p.level = 1

# SLIDE 8: TIEMPO
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = '📊 DISTRIBUCIÓN DEL TIEMPO'
slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = PptRGBColor(0, 102, 204)
table = slide.shapes.add_table(8, 3, PptInches(1), PptInches(2), PptInches(8), PptInches(4)).table
table.cell(0, 0).text = 'Sección'
table.cell(0, 1).text = 'Duración'
table.cell(0, 2).text = 'Acumulado'
for i, (sec, dur, acum) in enumerate(data, start=1):
    table.cell(i, 0).text = sec
    table.cell(i, 1).text = dur
    table.cell(i, 2).text = acum

# SLIDE 9: CIERRE
slide = prs.slides.add_slide(prs.slide_layouts[6])
txBox = slide.shapes.add_textbox(PptInches(1), PptInches(2.5), PptInches(8), PptInches(2))
tf = txBox.text_frame
tf.text = '🎉 ¡Éxito con tu grabación!'
tf.paragraphs[0].font.size = PptPt(54)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = PptRGBColor(0, 102, 204)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = '👨‍💻 Fuad Oñate\n🎓 Duoc UC'
p.font.size = PptPt(28)
p.alignment = PP_ALIGN.CENTER

ppt_file = 'GUION_VIDEO_FINAL_Fuad_Onate.pptx'
prs.save(ppt_file)
print(f'✅ PowerPoint guardado: {ppt_file}')

print('\n' + '='*60)
print('🎉 ¡ARCHIVOS FINALES CREADOS EXITOSAMENTE!')
print('='*60)
print(f'📄 Word: {word_file}')
print(f'📊 PowerPoint: {ppt_file}')
print('\n✨ Con tildes corregidas e íconos de programación')
